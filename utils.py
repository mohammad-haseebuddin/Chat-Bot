import speech_recognition as sr
from gtts import gTTS
import tempfile
import os
import webbrowser
from PIL import Image
import fitz
import time
import pandas as pd
import openpyxl
import pptx
import docx
from langchain.memory import StreamlitChatMessageHistory
import google.generativeai as genai
import streamlit as st

def speak_text(text):
    try:
        tts = gTTS(text=text, lang="en")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            tts.save(fp.name)
            audio_path = fp.name
        st.audio(audio_path, format="audio/mp3")
    except Exception as e:
        st.error(f"TTS Error: {e}")

def speech_to_text_from_mic():
    try:
        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            audio = recognizer.listen(source)
            text = recognizer.recognize_google(audio)
            return text
    except Exception as e:
        return f"Speech-to-Text Error: {e}"

def extract_text_from_file(file):
    try:
        file_extension = os.path.splitext(file.name)[1].lower()
        
        if file_extension == ".pdf":
            doc = fitz.open(stream=file.read(), filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        elif file_extension in [".txt", ".py", ".c", ".cpp", ".java", ".js", ".html", ".css"]:
            return file.read().decode("utf-8")
        elif file_extension in [".xlsx", ".xls"]:
            df = pd.read_excel(file)
            return df.to_string()
        elif file_extension in [".pptx", ".ppt"]:
            presentation = pptx.Presentation(file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        text += shape.text_frame.text
            return text
        elif file_extension == ".docx":
            doc = docx.Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"File processing error: {e}"

def handle_command(user_input):
    commands = {
        "youtube": "https://www.youtube.com",
        "google": "https://www.google.com",
        "github": "https://github.com",
    }

    user_input = user_input.lower()

    if "open" in user_input:
        for key, url in commands.items():
            if key in user_input:
                webbrowser.open(url)
                return f"Opening {key}..."

        if "notepad" in user_input:
            os.startfile("notepad.exe")
            return "Opening Notepad..."

        words = user_input.split()
        for word in words:
            if "." in word:
                if not word.startswith("http"):
                    word = "https://" + word
                webbrowser.open(word)
                return f"Opening {word}..."

    return None

def ask_gemini(prompt, model, image=None, chat_history=None, retries=3):
    full_prompt = ""
    if chat_history:
        for msg in chat_history:
            if hasattr(msg, 'type') and hasattr(msg, 'content'):
                if msg.type == "human":
                    full_prompt += f"User: {msg.content}\n"
                elif msg.type == "ai":
                    full_prompt += f"Assistant: {msg.content}\n"
    full_prompt += f"User: {prompt}\nAssistant:"

    for i in range(retries):
        try:
            if image:
                response = model.generate_content([full_prompt, image])
            else:
                response = model.generate_content(full_prompt)

            if response.candidates:
                candidate = response.candidates[0]
                if hasattr(candidate, "content") and candidate.content.parts:
                    texts = [part.text for part in candidate.content.parts if hasattr(part, "text")]
                    if texts:
                        return " ".join(texts)
            return "Sorry, I couldn’t generate a response."
        except Exception as e:
            if i < retries - 1:
                time.sleep(2**(i+1))
            else:
                return f"API Error: {e}"
    return "Failed to get a response after multiple retries."